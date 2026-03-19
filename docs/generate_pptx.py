"""Generate presentation slides for AI CallCenter CRM System."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
PRIMARY = RGBColor(0x1a, 0x56, 0xdb)
PRIMARY_DARK = RGBColor(0x1e, 0x40, 0xaf)
DARK = RGBColor(0x1f, 0x29, 0x37)
GRAY = RGBColor(0x6b, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0x9c, 0xa3, 0xaf)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG = RGBColor(0xf3, 0xf4, 0xf6)
ACCENT_GREEN = RGBColor(0x05, 0x96, 0x69)
ACCENT_ORANGE = RGBColor(0xd9, 0x77, 0x06)
ACCENT_RED = RGBColor(0xdc, 0x26, 0x26)
TABLE_HEADER_BG = RGBColor(0x1e, 0x40, 0xaf)
TABLE_ALT_BG = RGBColor(0xf0, 0xf4, 0xff)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Meiryo"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_screenshot_placeholder(slide, left, top, width, height, label=""):
    """Add a placeholder box for screenshot."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xfa, 0xfc)
    shape.line.color.rgb = RGBColor(0xd1, 0xd5, 0xdb)
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label if label else "Screenshot"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Meiryo"
    # vertical center
    tf.paragraphs[0].space_before = Pt(height.inches * 72 / 2 - 10)
    return shape


def add_table(slide, left, top, width, rows_data, col_widths_ratio=None):
    """Add a styled table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    if n_rows == 0 or n_cols == 0:
        return

    row_height = Inches(0.35)
    table_height = row_height * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height).table

    if col_widths_ratio:
        total = sum(col_widths_ratio)
        for i, ratio in enumerate(col_widths_ratio):
            table.columns[i].width = int(width * ratio / total)

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "Meiryo"
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table


def add_bullet_list(slide, left, top, width, items, font_size=12):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0
        # bullet
        p.text = "  " + item
    return txBox


def slide_title_bar(slide, title, subtitle=""):
    """Add a consistent title bar to content slides."""
    # Blue bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_DARK
    bar.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
                 title, font_size=24, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.35),
                     subtitle, font_size=12, color=RGBColor(0xbf, 0xdb, 0xfe))


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ===== SLIDE 1: TITLE =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PRIMARY_DARK)

    # Decorative circles
    for cx, cy, size, alpha in [(11.5, 1, 3, 0.08), (1, 6, 2.5, 0.06), (12, 6.5, 1.5, 0.1)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.fill.fore_color.brightness = 0.85
        c.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
                 "AI CallCenter CRM System", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
                 "機能紹介プレゼンテーション", font_size=24, color=RGBColor(0xbf, 0xdb, 0xfe), alignment=PP_ALIGN.CENTER)

    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.0), Inches(3.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x60, 0xa5, 0xfa)
    line.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                 "2026-03-18", font_size=16, color=RGBColor(0x93, 0xc5, 0xfd), alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 2: OVERVIEW =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "システム概要", "AI CallCenter CRM System の全体像")

    # Tech stack cards
    techs = [
        ("Frontend", "Next.js (React)", PRIMARY),
        ("Backend", "Express.js (Node.js)", ACCENT_GREEN),
        ("Database", "MySQL (Railway)", ACCENT_ORANGE),
        ("AI", "OpenAI GPT", ACCENT_RED),
        ("Hosting", "Railway", GRAY),
        ("外部連携", "Google Sheets API", PRIMARY),
    ]
    for i, (label, tech, color) in enumerate(techs):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4)
        y = Inches(1.5 + row * 1.4)
        card = add_shape(slide, x, y, Inches(3.5), Inches(1.1), WHITE, RGBColor(0xe5, 0xe7, 0xeb))
        # color bar on left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(3), Inches(0.35),
                     label, font_size=11, color=GRAY)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(3), Inches(0.4),
                     tech, font_size=16, color=DARK, bold=True)

    # Auth info
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.4),
                 "認証方式", font_size=16, color=DARK, bold=True)
    auth_items = [
        "  オペレーター: 名前選択でログイン（パスワードなし）",
        "  管理者/マネージャー/営業: メールアドレス + パスワード",
        "  JWT トークンによるセッション管理",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(auth_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(4)

    add_screenshot_placeholder(slide, Inches(7), Inches(4.4), Inches(5.5), Inches(2.7), "ログイン画面")

    # ===== SLIDE 3: USER ROLES =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "ユーザー権限", "4つのロールとアクセス制御")

    add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), [
        ["権限", "説明", "アクセス可能な機能"],
        ["operator", "オペレーター", "ダッシュボード、架電画面、リコール、案件管理、AI評価、架電結果、架電リスト、メッセージ"],
        ["manager", "マネージャー", "上記 + パフォーマンス閲覧、スクリプト管理、CPA分析、架電ログ管理"],
        ["admin", "管理者", "全機能（ユーザー管理含む）"],
        ["sales", "営業", "ダッシュボード、案件管理（閲覧中心）、メッセージ"],
    ], [1.5, 2, 8])

    # Visual role cards
    roles = [
        ("Operator", "架電業務の実行", PRIMARY),
        ("Manager", "チーム管理・分析", ACCENT_GREEN),
        ("Admin", "全機能アクセス", ACCENT_ORANGE),
        ("Sales", "案件管理・閲覧", ACCENT_RED),
    ]
    for i, (role, desc, color) in enumerate(roles):
        x = Inches(0.8 + i * 3.1)
        y = Inches(4.0)
        card = add_shape(slide, x, y, Inches(2.8), Inches(1.2), WHITE, color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.4),
                     role, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.4), Inches(0.35),
                     desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 4: DASHBOARD =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "ダッシュボード", "画面: /  |  権限: 全ロール")

    # KPI cards mockup
    kpis = [
        ("稼働時間", "7.0時間"),
        ("コール数", "46件"),
        ("リコール獲得", "7件"),
        ("リコール消化", "2件"),
        ("有効接続", "10件"),
    ]
    for i, (label, val) in enumerate(kpis):
        x = Inches(0.5 + i * 2.45)
        y = Inches(1.4)
        card = add_shape(slide, x, y, Inches(2.2), Inches(1.0), WHITE, RGBColor(0xe5, 0xe7, 0xeb))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(1.9), Inches(0.3),
                     label, font_size=10, color=GRAY)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(1.9), Inches(0.45),
                     val, font_size=20, color=DARK, bold=True)

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(5), Inches(0.35),
                 "期間タブ: 日別 / 週別 / 月別 / 累計  -- 全KPI・グラフが連動", font_size=11, color=PRIMARY, bold=True)

    add_screenshot_placeholder(slide, Inches(0.5), Inches(3.2), Inches(5.8), Inches(3.8), "ダッシュボード画面")

    # Feature list on right
    features = [
        "KPI カード (7指標)",
        "  稼働時間、コール数、リコール獲得/消化",
        "  有効接続、担当接続、案件獲得",
        "",
        "グラフ (3種)",
        "  時間帯別コール数 (棒グラフ)",
        "  業種別案件化率 (円グラフ)",
        "  時間帯x業種別 接続率 (テーブル)",
        "",
        "コピー機能",
        "  コールデータ / 日報コピー",
        "",
        "AI総合分析 (管理者/マネージャー)",
        "  チーム / 個人のパフォーマンスをAIが分析",
    ]
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(3.2), Inches(6), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item == "":
            p.text = ""
            p.space_after = Pt(2)
            continue
        is_header = not item.startswith("  ")
        p.text = item
        p.font.size = Pt(11 if is_header else 10)
        p.font.color.rgb = DARK if is_header else GRAY
        p.font.bold = is_header
        p.font.name = "Meiryo"
        p.space_after = Pt(2)

    # ===== SLIDE 5: CALL SCREEN =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "架電画面", "画面: /call  |  権限: operator, manager, admin")

    add_screenshot_placeholder(slide, Inches(0.5), Inches(1.4), Inches(7), Inches(5.5), "架電画面")

    # Features on right
    sections = [
        ("架電リスト", [
            "最大10件の架電対象を表示",
            "4つのピックアップモード:",
            "  自動 / 業種別 / 自作 / 特別",
        ]),
        ("企業情報", [
            "会社名クリックでGoogle検索",
            "電話番号、業種、地域、住所",
            "過去の架電履歴表示",
        ]),
        ("結果入力", [
            "6種の結果コード",
            "メモ、接続フラグ",
            "リコール日時設定",
        ]),
        ("排他ロック", [
            "企業単位で5分間ロック",
            "他オペレーターとの重複防止",
        ]),
    ]
    y_pos = Inches(1.4)
    for title, items in sections:
        add_text_box(slide, Inches(7.8), y_pos, Inches(5), Inches(0.3),
                     title, font_size=13, color=PRIMARY, bold=True)
        y_pos += Inches(0.35)
        for item in items:
            add_text_box(slide, Inches(7.8), y_pos, Inches(5), Inches(0.25),
                         item, font_size=10, color=DARK if not item.startswith("  ") else GRAY)
            y_pos += Inches(0.25)
        y_pos += Inches(0.15)

    # ===== SLIDE 6: CALL RESULTS =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "架電結果 / 結果コード", "画面: /call-results  |  権限: operator, manager, admin")

    add_table(slide, Inches(0.8), Inches(1.5), Inches(5.5), [
        ["結果コード", "説明"],
        ["NO_ANSWER", "不通"],
        ["NG", "不在/拒否"],
        ["RECALL", "リコール（日時指定）"],
        ["INTERESTED", "興味あり"],
        ["PROJECT", "案件化"],
        ["SKIP", "スキップ（集計対象外）"],
    ], [2, 3.5])

    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.35),
                 "架電結果画面の主な機能", font_size=14, color=DARK, bold=True)
    result_features = [
        "  日別/範囲/全件の表示切替",
        "  フィルタ: 結果コード、オペレーター、検索",
        "  インライン編集: 結果コード、メモ、接続フラグ",
        "  通話時間の自動計算・表示",
        "  文字起こし展開表示",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(5.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(result_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    add_screenshot_placeholder(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), "架電結果一覧")

    # ===== SLIDE 7: RECALL =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "リコール管理", "画面: /recalls  |  権限: operator, manager, admin")

    add_screenshot_placeholder(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(4), "リコール管理画面")

    add_table(slide, Inches(7.5), Inches(1.5), Inches(5.3), [
        ["タブ", "対象"],
        ["今日", "本日期限のリコール"],
        ["明日", "翌日期限のリコール"],
        ["期限超過", "期限切れの未完了リコール"],
        ["その他", "上記以外の予定リコール"],
    ], [1.5, 3.8])

    add_text_box(slide, Inches(7.5), Inches(3.8), Inches(5.3), Inches(0.35),
                 "操作", font_size=14, color=DARK, bold=True)
    recall_ops = [
        "  完了: リコールを消化済みにする",
        "  キャンセル: リコールを取消す",
        "  通話文字起こしの閲覧",
        "  クリックで架電画面へ遷移",
    ]
    txBox = slide.shapes.add_textbox(Inches(7.5), Inches(4.2), Inches(5.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(recall_ops):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    # ===== SLIDE 8: PROJECTS =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "案件管理", "画面: /projects  |  権限: 全ロール")

    add_screenshot_placeholder(slide, Inches(0.5), Inches(1.4), Inches(6.5), Inches(3), "案件一覧画面")

    add_table(slide, Inches(7.5), Inches(1.4), Inches(5.3), [
        ["ステータス", "説明"],
        ["BOSHUCHU", "募集中"],
        ["SHORUI_CHU", "書類選考中"],
        ["MENSETSU_KAKUTEI", "面接確定"],
        ["KEKKA_MACHI", "結果待ち"],
        ["NAITEI", "内定"],
        ["NAITEI_TORIKESHI", "内定取消"],
        ["FUGOKAKU", "不合格"],
        ["LOST", "失注"],
        ["BARASHI", "バラシ"],
    ], [2, 3.3])

    # Key features
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(6), Inches(0.35),
                 "主な機能", font_size=14, color=DARK, bold=True)
    proj_features = [
        "  インライン ステータス変更（プルダウンで直接変更）",
        "  ステータスフィルタ、日付範囲、自分の案件フィルタ",
        "  案件詳細: 求人番号、面接日、面接種別、書類選考、メモ",
        "  過去の架電履歴タイムライン",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(proj_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    # ===== SLIDE 9: HIRE INFO =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "内定者情報", "案件管理 - 内定時の追加機能")

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.35),
                 'ステータスを「内定」に変更すると自動でモーダルが表示される（一覧画面・詳細画面両対応）', font_size=13, color=PRIMARY)

    add_table(slide, Inches(0.8), Inches(2.2), Inches(5.5), [
        ["入力項目", "説明", "備考"],
        ["登録番号", "候補者の登録番号", "例: AB1234"],
        ["コース", "国内/転職/海外", "プルダウン選択"],
        ["初回入金", "初回入金額", "半角数字のみ"],
        ["見込売上", "見込売上額", "半角数字のみ"],
    ], [1.5, 2, 2])

    add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5), "内定者情報モーダル")

    # Cancel info
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.35),
                 "内定取消の仕組み", font_size=14, color=DARK, bold=True)
    cancel_items = [
        '  ステータス「内定取消」で全員の金額が0にリセット',
        "  個別取消: チェックボックスで1人ずつ取消可能",
        "  復活: チェックを外すと金額を再編集可能（可逆）",
        "  CPA分析では取消分を自動除外",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(5.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(cancel_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    # ===== SLIDE 10: AI EVAL =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "AI評価", "画面: /logs  |  権限: operator, manager, admin")

    # 5 dimension cards
    dims = [
        ("オープニング", "挨拶・自己紹介・用件提示"),
        ("明瞭さ", "話し方の分かりやすさ"),
        ("ヒアリング", "相手のニーズ把握力"),
        ("切り返し", "反論・断りへの対応力"),
        ("クロージング", "次アクションへの誘導力"),
    ]
    for i, (dim, desc) in enumerate(dims):
        x = Inches(0.5 + i * 2.45)
        y = Inches(1.4)
        card = add_shape(slide, x, y, Inches(2.2), Inches(1.1), WHITE, PRIMARY)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.9), Inches(0.35),
                     dim, font_size=13, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.9), Inches(0.4),
                     desc, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(2.7), Inches(6), Inches(0.3),
                 "各項目 100点満点で評価  |  15行未満の通話は自動スキップ", font_size=11, color=GRAY)

    add_screenshot_placeholder(slide, Inches(0.5), Inches(3.2), Inches(6), Inches(4), "AI評価画面")

    # Features on right
    eval_sections = [
        ("評価機能", [
            "個別評価: 1件ずつAI評価",
            "一括評価: 日次の全通話を評価",
            "日次サマリー: 平均スコアと傾向",
        ]),
        ("フィードバック", [
            "良かった点（箇条書き）",
            "改善点（箇条書き）",
            "次回アクション提案",
        ]),
        ("管理者向け", [
            "全オペレーターの評価一覧",
            "スクリプト提案の自動生成",
            "電話番号での横断検索",
        ]),
    ]
    y_pos = Inches(3.2)
    for title, items in eval_sections:
        add_text_box(slide, Inches(7), y_pos, Inches(5.5), Inches(0.3),
                     title, font_size=13, color=PRIMARY, bold=True)
        y_pos += Inches(0.35)
        for item in items:
            add_text_box(slide, Inches(7.3), y_pos, Inches(5.2), Inches(0.25),
                         "  " + item, font_size=10, color=DARK)
            y_pos += Inches(0.25)
        y_pos += Inches(0.15)

    # ===== SLIDE 11: CSV IMPORT =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "架電リスト管理（CSVインポート）", "画面: /csv-import  |  権限: operator, manager, admin")

    add_table(slide, Inches(0.8), Inches(1.5), Inches(6), [
        ["種別", "説明", "権限"],
        ["架電リスト", "通常の架電対象企業", "全員"],
        ["特別リスト", "特別ピックアップ用", "manager以上"],
        ["NGリスト", "架電除外企業", "manager以上"],
        ["既存案件", "案件化済み（除外用）", "manager以上"],
    ], [1.5, 3, 1.5])

    add_screenshot_placeholder(slide, Inches(7.3), Inches(1.5), Inches(5.5), Inches(3), "CSVインポート画面")

    imp_features = [
        ("インポート方法", "CSV / XLS / XLSX（ドラッグ&ドロップ可）、手動入力"),
        ("重複判定", "電話番号の一致で自動判定"),
        ("優先割当", "オペレーター + 猶予日数を指定（manager以上）"),
        ("企業一覧", "業種・地域フィルタ、ロック状態表示（30秒更新）"),
    ]
    y_pos = Inches(4.0)
    for title, desc in imp_features:
        add_text_box(slide, Inches(0.8), y_pos, Inches(2.5), Inches(0.25),
                     title, font_size=11, color=PRIMARY, bold=True)
        add_text_box(slide, Inches(3.3), y_pos, Inches(9.5), Inches(0.25),
                     desc, font_size=11, color=DARK)
        y_pos += Inches(0.4)

    # ===== SLIDE 12: CPA ANALYTICS =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "CPA/案件質分析", "画面: /admin/analytics  |  権限: manager, admin")

    add_screenshot_placeholder(slide, Inches(0.5), Inches(1.4), Inches(7), Inches(5.5), "CPA分析画面")

    # CPA metrics
    add_text_box(slide, Inches(7.8), Inches(1.4), Inches(5), Inches(0.35),
                 "CPA指標 (13項目)", font_size=14, color=PRIMARY, bold=True)
    cpa_items = [
        "コスト = 出勤記録 x 時給1,500円",
        "コール数 / 案件化率 / 案件数",
        "案件CPA / 面接CPA",
        "内定 / 不合格 / バラシ・失注",
        "初回入金 / 見込売上 / ROAS",
    ]
    y_pos = Inches(1.9)
    for item in cpa_items:
        add_text_box(slide, Inches(8), y_pos, Inches(4.8), Inches(0.25),
                     "  " + item, font_size=10, color=DARK)
        y_pos += Inches(0.27)

    add_text_box(slide, Inches(7.8), Inches(3.5), Inches(5), Inches(0.35),
                 "案件質指標 (9項目)", font_size=14, color=PRIMARY, bold=True)
    quality_items = [
        "案件数 / 失注 / 連絡待ち",
        "面接日確定 / 面接実施",
        "バラシ / オンライン面接",
        "書類選考無し / 書類選考落ち",
    ]
    y_pos = Inches(3.9)
    for item in quality_items:
        add_text_box(slide, Inches(8), y_pos, Inches(4.8), Inches(0.25),
                     "  " + item, font_size=10, color=DARK)
        y_pos += Inches(0.27)

    add_text_box(slide, Inches(7.8), Inches(5.2), Inches(5), Inches(0.35),
                 "表示形式", font_size=14, color=PRIMARY, bold=True)
    display_items = [
        "全オペレーター比較テーブル",
        "期間: 月別 / 週別（全週一覧）/ 累計",
        "コスト: CSV/PDFインポート",
    ]
    y_pos = Inches(5.6)
    for item in display_items:
        add_text_box(slide, Inches(8), y_pos, Inches(4.8), Inches(0.25),
                     "  " + item, font_size=10, color=DARK)
        y_pos += Inches(0.27)

    # ===== SLIDE 13: SCRIPTS =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "スクリプト管理 / ユーザー管理", "管理者向け機能")

    # Scripts section
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.35),
                 "スクリプト管理  /admin/scripts", font_size=15, color=DARK, bold=True)

    add_table(slide, Inches(0.8), Inches(2.0), Inches(5.5), [
        ["種別", "説明"],
        ["切り返しトーク", "断り文句への対応スクリプト"],
        ["Q&A", "よくある質問と回答"],
    ], [2, 3.5])

    script_flow = [
        "1. 管理者/マネージャーがスクリプトを作成",
        "2. 承認/却下の管理",
        "3. 承認済みがオペレーターの架電画面に表示",
        "4. 業種フィルタリング対応",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(script_flow):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "  " + item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    # User management section
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.35),
                 "ユーザー管理  /admin/users", font_size=15, color=DARK, bold=True)
    user_features = [
        "  ユーザーの作成/編集/削除",
        "  ロール割当（operator / manager / admin / sales）",
        "  アクティブ/非アクティブ切替",
        "  一覧表示（名前、メール、ロール、作成日）",
    ]
    txBox = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(user_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    # Message section
    add_text_box(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(0.35),
                 "メッセージ機能  /messages", font_size=15, color=DARK, bold=True)
    msg_features = [
        "  ユーザーが機能要望やフィードバックを送信",
        "  管理者が全メッセージを閲覧・返信",
        "  メッセージ一覧と返信履歴の表示",
    ]
    txBox = slide.shapes.add_textbox(Inches(7), Inches(4.2), Inches(5.5), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(msg_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.font.name = "Meiryo"
        p.space_after = Pt(3)

    add_screenshot_placeholder(slide, Inches(0.5), Inches(5.0), Inches(5.8), Inches(2), "スクリプト管理画面")
    add_screenshot_placeholder(slide, Inches(7), Inches(5.5), Inches(5.5), Inches(1.5), "ユーザー管理画面")

    # ===== SLIDE 14: DB SCHEMA =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    slide_title_bar(slide, "データベース構成", "主要テーブル一覧")

    add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), [
        ["テーブル", "説明", "主な用途"],
        ["users", "ユーザー情報（ロール、アクティブ状態）", "認証・権限管理"],
        ["companies", "企業情報（業種、地域、ロック状態）", "架電対象管理"],
        ["calls", "架電記録（結果、メモ、接続フラグ）", "架電履歴・KPI集計"],
        ["projects", "案件情報（ステータス、面接情報）", "案件管理・進捗追跡"],
        ["project_hires", "内定者情報（登録番号、コース、入金）", "内定管理・CPA算出"],
        ["recall_tasks", "リコール予定（日時、ステータス）", "リコール管理"],
        ["evaluations", "AI評価結果（5次元スコア）", "品質管理・フィードバック"],
        ["scripts", "スクリプト（種別、承認状態）", "トーク品質向上"],
        ["cost_records", "出勤記録（CSV/PDFインポート）", "CPA算出用コストデータ"],
        ["exclusion_lists", "NG/既存案件リスト", "架電除外管理"],
        ["feature_requests", "ユーザーからの要望", "フィードバック管理"],
    ], [2, 4.5, 5.2])

    # ===== SLIDE 15: END =====
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PRIMARY_DARK)

    for cx, cy, size in [(11, 1.5, 2.5), (1.5, 5.5, 2), (12, 6, 1.5)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.fill.fore_color.brightness = 0.85
        c.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                 "AI CallCenter CRM System", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.7), Inches(3.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x60, 0xa5, 0xfa)
    line.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
                 "ご清聴ありがとうございました", font_size=22, color=RGBColor(0xbf, 0xdb, 0xfe), alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
    prs.save(output_path)
    print(f"PPTX generated: {output_path}")


if __name__ == "__main__":
    build_pptx()
